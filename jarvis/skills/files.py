"""File operations - create, delete, search files."""
import os
import subprocess
from pathlib import Path
from datetime import datetime

def handle(query: str) -> str:
    """Handle file operations."""
    q = query.lower()
    
    # Create
    if any(kw in q for kw in ["create", "make", "new"]):
        return create_file(q)
    
    # Delete
    if any(kw in q for kw in ["delete", "remove"]):
        return delete_file(q)
    
    # Search (New)
    if any(kw in q for kw in ["find", "search", "locate", "where"]):
        return search_files(q)
        
    # List
    if "list" in q:
        return list_files()
    
    return None

def search_files(query: str) -> str:
    """Smart file search in user directories."""
    q = query.lower()
    
    # 1. Parse Filters
    extensions = []
    if "pdf" in q: extensions.append(".pdf")
    if "word" in q or "doc" in q: extensions.extend([".docx", ".doc"])
    if "image" in q or "photo" in q or "picture" in q: extensions.extend([".jpg", ".png", ".jpeg"])
    if "presentation" in q or "ppt" in q: extensions.extend([".pptx", ".ppt"])
    if "text" in q or "txt" in q: extensions.append(".txt")
    if "excel" in q or "sheet" in q: extensions.extend([".xlsx", ".xls", ".csv"])
    
    max_days = 3650
    if "today" in q: max_days = 1
    elif "yesterday" in q: max_days = 2
    elif "this week" in q or "recent" in q: max_days = 7
    elif "last week" in q: max_days = 14
    elif "month" in q: max_days = 30
    
    triggers = ["find", "search", "locate", "where is", "where are", "look for", "show me", "the", "file", "files", "list", "all", "created", "from", "in"]
    search_term = q
    for t in triggers:
        search_term = search_term.replace(f" {t} ", " ") # Replace match with spaces
        if search_term.startswith(f"{t} "): search_term = search_term[len(t)+1:]
    
    filter_words = ["pdf", "word", "doc", "image", "photo", "today", "yesterday", "recent", "week", "month", "text", "txt", "excel", "presentation", "ppt"]
    for w in filter_words:
        search_term = search_term.replace(w, "")
        
    search_term = search_term.strip()
    
    user_home = Path.home()
    locations = [
        user_home / "Desktop",
        user_home / "Documents",
        user_home / "Downloads"
    ]
    
    results = []
    limit = 5
    now = datetime.now()
    
    print(f"🔍 Searching for '{search_term}' (Ext: {extensions}, Days: {max_days})...")
    
    for loc in locations:
        if not loc.exists(): continue
        try:
            for root, _, files in os.walk(loc):
                if len(results) >= limit: break
                if any(x in root.lower() for x in ["node_modules", ".git", "appdata", "library"]): continue
                    
                for file in files:
                    if len(results) >= limit: break
                    if extensions and not any(file.lower().endswith(ext) for ext in extensions): continue
                    if search_term and search_term not in file.lower(): continue
                    
                    path = Path(root) / file
                    try:
                        mtime = datetime.fromtimestamp(path.stat().st_mtime)
                        age_days = (now - mtime).days
                        if age_days > max_days: continue
                        results.append(path)
                    except: pass
        except: pass
            
    if results:
        results.sort(key=lambda x: x.stat().st_mtime, reverse=True)
        response = "Found these files:\n"
        for p in results:
            response += f"- {p.name} ({p.parent.name})\n"
        return response.strip()
        
    return f"I couldn't find any matching files for '{search_term}'."

def create_file(query: str) -> str:
    """Create a file."""
    downloads = Path.home() / "Downloads"
    downloads.mkdir(exist_ok=True)
    
    name = query.split("create")[-1].split("make")[-1].split("new")[-1].strip()
    if not name:
        name = f"file_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
    
    if "." not in name:
        if "word" in query or "doc" in query:
            name += ".docx"
            return create_word(downloads / name)
        elif "pdf" in query:
            name += ".pdf"
            return create_pdf(downloads / name)
        elif "ppt" in query or "presentation" in query:
            name += ".pptx"
            return create_ppt(downloads / name)
        else:
            name += ".txt"
    
    filepath = downloads / name
    filepath.touch()
    try:
        subprocess.Popen(["notepad.exe", str(filepath)])
    except: pass
    
    return f"Created {name} in Downloads"

def create_word(path: Path) -> str:
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("Document", 0)
        doc.add_paragraph("Created by JARVIS")
        doc.save(str(path))
        subprocess.Popen(["start", str(path)], shell=True)
        return f"Created Word document: {path.name}"
    except: return "Install python-docx: pip install python-docx"

def create_pdf(path: Path) -> str:
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(str(path), pagesize=letter)
        c.drawString(100, 750, "Document Created by JARVIS")
        c.save()
        subprocess.Popen(["start", str(path)], shell=True)
        return f"Created PDF: {path.name}"
    except: return "Install reportlab: pip install reportlab"

def create_ppt(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Presentation"
        slide.placeholders[1].text = "Created by JARVIS"
        prs.save(str(path))
        subprocess.Popen(["start", str(path)], shell=True)
        return f"Created presentation: {path.name}"
    except: return "Install python-pptx: pip install python-pptx"

def delete_file(query: str) -> str:
    if "confirm" not in query:
        return "Say 'delete file confirm' to delete"
    
    downloads = Path.home() / "Downloads"
    files = list(downloads.glob("*"))
    if files:
        files.sort(key=lambda x: x.stat().st_mtime, reverse=True)
        files[0].unlink()
        return f"Deleted {files[0].name}"
    return "No files found"

def list_files() -> str:
    downloads = Path.home() / "Downloads"
    files = [f.name for f in downloads.glob("*") if f.is_file()][:10]
    if files:
        return "Files: " + ", ".join(files)
    return "No files in Downloads"
